"""
PPTX generation service.
Builds a professional white + deep-purple deck from a structured JSON slide outline.
"""

from __future__ import annotations


def outline_to_pptx(outline: dict, output_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # ── Palette ────────────────────────────────────────────────────
    PURPLE       = RGBColor(0x6D, 0x28, 0xD9)
    PURPLE_DARK  = RGBColor(0x4C, 0x1D, 0x95)
    PURPLE_LIGHT = RGBColor(0xED, 0xE9, 0xFE)
    PURPLE_ACCENT= RGBColor(0xA7, 0x8B, 0xFA)
    TEXT_DARK    = RGBColor(0x1F, 0x29, 0x37)
    TEXT_MUTED   = RGBColor(0x6B, 0x72, 0x80)
    TEXT_FAINT   = RGBColor(0x9C, 0xA3, 0xAF)
    WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER      = RGBColor(0xE5, 0xE7, 0xEB)
    SLATE_50     = RGBColor(0xF8, 0xFA, 0xFC)
    SLATE_100    = RGBColor(0xF1, 0xF5, 0xF9)

    # Severity palette
    SEV_CRITICAL = RGBColor(0xDC, 0x26, 0x26)
    SEV_HIGH     = RGBColor(0xEA, 0x58, 0x0C)
    SEV_MEDIUM   = RGBColor(0xCA, 0x8A, 0x04)
    SEV_LOW      = RGBColor(0x15, 0x80, 0x3D)
    SEV_INFO     = RGBColor(0x0E, 0x7A, 0xBF)

    SEV_CRIT_LIGHT = RGBColor(0xFE, 0xF2, 0xF2)
    SEV_HIGH_LIGHT = RGBColor(0xFF, 0xF7, 0xED)
    SEV_MED_LIGHT  = RGBColor(0xFF, 0xFB, 0xEB)
    SEV_LOW_LIGHT  = RGBColor(0xF0, 0xFD, 0xF4)
    SEV_INFO_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)

    SEV_MAP = {
        'critical': (SEV_CRITICAL, SEV_CRIT_LIGHT),
        'high':     (SEV_HIGH,     SEV_HIGH_LIGHT),
        'medium':   (SEV_MEDIUM,   SEV_MED_LIGHT),
        'med':      (SEV_MEDIUM,   SEV_MED_LIGHT),
        'low':      (SEV_LOW,      SEV_LOW_LIGHT),
        'info':     (SEV_INFO,     SEV_INFO_LIGHT),
        'neutral':  (PURPLE,       PURPLE_LIGHT),
        'positive': (SEV_LOW,      SEV_LOW_LIGHT),
        'negative': (SEV_CRITICAL, SEV_CRIT_LIGHT),
    }

    def sev_colors(key: str):
        return SEV_MAP.get((key or '').lower(), (PURPLE, PURPLE_LIGHT))

    FONT_SANS = "Aptos"
    SLIDE_W   = Inches(13.333)
    SLIDE_H   = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK_LAYOUT = prs.slide_layouts[6]

    meta           = outline.get("meta", {}) if isinstance(outline, dict) else {}
    deck_title     = meta.get("title", "Report")
    classification = meta.get("classification", "TLP:AMBER+STRICT")
    date_str       = meta.get("date", "")
    author         = meta.get("author", "SkillCTI")

    def add_blank_slide():
        return prs.slides.add_slide(BLANK_LAYOUT)

    def set_text(tf, text, *, size=12, bold=False, color=TEXT_DARK,
                 font=FONT_SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for p in list(tf.paragraphs)[1:]:
            tf._txBody.remove(p._p)
        if not tf.paragraphs:
            tf.add_paragraph()
        p = tf.paragraphs[0]
        p.alignment = align
        pPr = p._p
        for r in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
            pPr.remove(r)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def add_textbox(slide, left, top, width, height, text, **kw):
        tb = slide.shapes.add_textbox(left, top, width, height)
        set_text(tb.text_frame, text, **kw)
        return tb

    def add_rect(slide, left, top, width, height, fill, line=None, line_width=0.75):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_width)
        shp.shadow.inherit = False
        return shp

    def add_footer(slide, page_num, total):
        add_rect(slide, Emu(0), Inches(7.22), SLIDE_W, Inches(0.04), PURPLE)
        footer_text = f"{classification}  ·  {deck_title}  ·  {author}  ·  {page_num} / {total}"
        add_textbox(
            slide, Inches(0.35), Inches(7.26), Inches(12.6), Inches(0.24),
            footer_text, size=7.5, color=TEXT_FAINT, font=FONT_SANS,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

    def add_slide_header(slide, title, subtitle=None, *, accent_color=PURPLE):
        """Draw header; returns the Y coordinate where content should start."""
        add_rect(slide, Inches(0.4), Inches(0.38), Inches(0.07), Inches(0.72), accent_color)
        add_textbox(slide, Inches(0.65), Inches(0.32), Inches(12.3), Inches(0.82),
                    title, size=24, bold=True, color=TEXT_DARK, font=FONT_SANS)
        if subtitle:
            add_textbox(slide, Inches(0.65), Inches(1.12), Inches(12.3), Inches(0.40),
                        subtitle, size=11, color=TEXT_MUTED, font=FONT_SANS)
        div_top = Inches(1.60) if subtitle else Inches(1.28)
        add_rect(slide, Inches(0.4), div_top, Inches(12.53), Emu(8000), DIVIDER)
        return int(div_top) + int(Inches(0.18))

    def style_run(run, text, *, size=14, bold=False, color=TEXT_DARK, font=FONT_SANS):
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def fill_bullet_paragraph(p, item, size, line_spacing, bullet_color):
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        pPr = p._p
        for r in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
            pPr.remove(r)
        style_run(p.add_run(), "▸  ", size=size - 1, bold=True, color=bullet_color)
        text = str(item)
        if ":" in text and text.index(":") < 50:
            head, tail = text.split(":", 1)
            style_run(p.add_run(), head + ":", size=size, bold=True, color=TEXT_DARK)
            style_run(p.add_run(), tail, size=size, color=TEXT_DARK)
        else:
            style_run(p.add_run(), text, size=size, color=TEXT_DARK)

    def add_bullet_list(slide, bullets, left, top, width, height,
                        size=14, line_spacing=1.45, bullet_color=PURPLE):
        if not bullets:
            return
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            fill_bullet_paragraph(p, item, size, line_spacing, bullet_color)

    # ── Slide Renderers ────────────────────────────────────────────

    def render_title(slide, s):
        # Purple left edge panel
        add_rect(slide, Emu(0), Emu(0), Inches(0.65), SLIDE_H, PURPLE)

        # Subtle horizontal stripe pattern on white area
        for i in range(9):
            add_rect(slide, Inches(0.85), Inches(0.85 * i + 0.1),
                     Inches(12.0), Emu(6000), DIVIDER)

        # Classification pill (top right)
        pill_w, pill_h = Inches(2.5), Inches(0.40)
        add_rect(slide, Inches(10.6), Inches(0.38), pill_w, pill_h, WHITE, line=PURPLE)
        add_textbox(slide, Inches(10.6), Inches(0.38), pill_w, pill_h,
                    classification, size=9.5, bold=True, color=PURPLE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Main title
        add_textbox(slide, Inches(1.0), Inches(2.1), Inches(11.8), Inches(2.0),
                    s.get("title", deck_title),
                    size=42, bold=True, color=PURPLE_DARK, font=FONT_SANS)

        # Accent underline
        add_rect(slide, Inches(1.0), Inches(4.15), Inches(1.4), Inches(0.12), PURPLE)

        sub = s.get("subtitle") or meta.get("subtitle", "")
        if sub:
            add_textbox(slide, Inches(1.0), Inches(4.4), Inches(11.5), Inches(0.75),
                        sub, size=17, color=TEXT_MUTED, font=FONT_SANS)

        client_line = meta.get("client", "")
        bottom = "  ·  ".join(x for x in [client_line, date_str, author] if x)
        if bottom:
            add_rect(slide, Inches(1.0), Inches(6.42), Inches(11.5), Emu(8000), DIVIDER)
            add_textbox(slide, Inches(1.0), Inches(6.55), Inches(11.5), Inches(0.45),
                        bottom, size=10.5, color=TEXT_FAINT, font=FONT_SANS)

    def render_section(slide, s):
        # Dark purple left panel (~40% width)
        panel_w = Inches(5.2)
        add_rect(slide, Emu(0), Emu(0), panel_w, SLIDE_H, PURPLE_DARK)
        # Light purple right panel
        add_rect(slide, panel_w, Emu(0), int(SLIDE_W) - int(panel_w), SLIDE_H, PURPLE_LIGHT)

        num = str(s.get("number", ""))
        if num:
            # Large faded number on the left panel
            add_textbox(slide, Inches(0.3), Inches(0.3), Inches(4.6), Inches(3.5),
                        num, size=160, bold=True,
                        color=RGBColor(0x5B, 0x21, 0xB6),  # slightly lighter purple
                        font=FONT_SANS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_textbox(slide, Inches(0.3), Inches(4.0), Inches(4.6), Inches(0.5),
                        "SECTION", size=9, bold=True, color=PURPLE_ACCENT,
                        font=FONT_SANS, align=PP_ALIGN.CENTER)

        # Title on right side
        add_textbox(slide, Inches(5.6), Inches(2.5), Inches(7.3), Inches(1.5),
                    s.get("title", "Section"),
                    size=36, bold=True, color=PURPLE_DARK, font=FONT_SANS)
        if s.get("subtitle"):
            add_textbox(slide, Inches(5.6), Inches(4.2), Inches(7.3), Inches(0.7),
                        s["subtitle"], size=15, color=TEXT_MUTED, font=FONT_SANS)

    def render_content(slide, s):
        content_y = add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        bullets = s.get("bullets", [])

        if bullets:
            # Left accent bar spanning the bullet area
            add_rect(slide, Inches(0.40), content_y, Inches(0.05), Inches(5.2), PURPLE_ACCENT)
            add_bullet_list(slide, bullets,
                            Inches(0.65), content_y + int(Pt(4)),
                            Inches(12.2), Inches(5.2),
                            size=14, line_spacing=1.5)

        if s.get("footnote"):
            add_textbox(slide, Inches(0.65), Inches(6.82), Inches(12.3), Inches(0.38),
                        s["footnote"], size=9, color=TEXT_FAINT, font=FONT_SANS)

    def render_stats(slide, s):
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        stats = s.get("stats", [])[:5]
        if not stats:
            return
        n = len(stats)
        gap = Inches(0.22)
        usable_w = Inches(12.13)
        card_w = int((int(usable_w) - int(gap) * (n - 1)) / n)
        left = Inches(0.6)
        top = Inches(2.15)
        card_h = Inches(3.3)

        for i, st in enumerate(stats):
            x = int(left) + i * (card_w + int(gap))
            fg, bg = sev_colors(st.get("severity", "") or st.get("color", ""))

            # Card background + border
            add_rect(slide, x, top, card_w, card_h, WHITE, line=DIVIDER, line_width=0.5)
            # Thick colored top bar
            add_rect(slide, x, top, card_w, Inches(0.12), fg)

            # Value
            value_text = str(st.get("value", ""))
            trend = st.get("trend", "")
            if trend == "up":
                value_text += " ▲"
            elif trend == "down":
                value_text += " ▼"

            add_textbox(slide, x, int(top) + int(Inches(0.3)), card_w, Inches(1.5),
                        value_text,
                        size=36, bold=True, color=PURPLE_DARK, font=FONT_SANS,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # Label
            add_textbox(slide, x, int(top) + int(Inches(1.95)), card_w, Inches(0.9),
                        str(st.get("label", "")),
                        size=11, color=TEXT_MUTED, font=FONT_SANS,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

            # Optional sublabel
            if st.get("sublabel"):
                add_textbox(slide, x, int(top) + int(Inches(2.75)), card_w, Inches(0.45),
                            str(st["sublabel"]),
                            size=9, color=TEXT_FAINT, font=FONT_SANS,
                            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    def render_table(slide, s):
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        headers = s.get("headers", [])
        rows = s.get("rows", [])
        if not headers or not rows:
            return
        n_cols = len(headers)
        n_rows = len(rows) + 1
        row_h = Inches(0.42)
        table_h = min(int(row_h) * n_rows + int(Inches(0.05)), int(Inches(5.0)))
        table = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.55), Inches(1.88),
            Inches(12.23), table_h,
        ).table

        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PURPLE_DARK
            set_text(cell.text_frame, str(h), size=10.5, bold=True,
                     color=WHITE, font=FONT_SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        for i, row in enumerate(rows, start=1):
            bg = WHITE if i % 2 else SLATE_50
            for j in range(n_cols):
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                value = str(row[j]) if j < len(row) else ""
                cell_color = TEXT_DARK
                if j == 0 and value.lower() in SEV_MAP:
                    cell_color, _ = sev_colors(value.lower())
                    value = value.upper()
                set_text(cell.text_frame, value, size=10, color=cell_color,
                         font=FONT_SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    def render_two_column(slide, s):
        add_slide_header(slide, s.get("title", ""))
        col_w = Inches(5.75)
        gap = Inches(0.48)
        left_x = Inches(0.55)
        right_x = int(left_x) + int(col_w) + int(gap)
        content_top = Inches(1.72)
        col_h = Inches(5.0)

        # Column background fills
        add_rect(slide, left_x, content_top, col_w, col_h, SLATE_50)
        add_rect(slide, right_x, content_top, col_w, col_h, SLATE_50)

        # Left header bar
        hdr_h = Inches(0.55)
        add_rect(slide, left_x, content_top, col_w, hdr_h, PURPLE)
        add_textbox(slide, int(left_x) + int(Inches(0.18)), int(content_top) + int(Inches(0.08)),
                    int(col_w) - int(Inches(0.3)), hdr_h,
                    s.get("left_title", ""), size=13, bold=True, color=WHITE,
                    font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)

        # Right header bar
        add_rect(slide, right_x, content_top, col_w, hdr_h, PURPLE_DARK)
        add_textbox(slide, right_x + int(Inches(0.18)), int(content_top) + int(Inches(0.08)),
                    int(col_w) - int(Inches(0.3)), hdr_h,
                    s.get("right_title", ""), size=13, bold=True, color=WHITE,
                    font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)

        bullet_top = int(content_top) + int(hdr_h) + int(Inches(0.15))
        add_bullet_list(slide, s.get("left_bullets", []),
                        int(left_x) + int(Inches(0.18)), bullet_top,
                        int(col_w) - int(Inches(0.3)), Inches(4.1), size=12.5)
        add_bullet_list(slide, s.get("right_bullets", []),
                        right_x + int(Inches(0.18)), bullet_top,
                        int(col_w) - int(Inches(0.3)), Inches(4.1), size=12.5)

    def render_quote(slide, s):
        add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, RGBColor(0xFB, 0xFA, 0xFE))
        add_textbox(slide, Inches(0.8), Inches(0.9), Inches(2.0), Inches(2.2),
                    "“", size=120, bold=True, color=PURPLE_ACCENT, font=FONT_SANS)
        add_textbox(slide, Inches(1.8), Inches(2.1), Inches(10.5), Inches(2.8),
                    s.get("quote", ""), size=26, color=TEXT_DARK, font=FONT_SANS)
        add_rect(slide, Inches(1.8), Inches(5.25), Inches(1.5), Inches(0.07), PURPLE)
        add_textbox(slide, Inches(1.8), Inches(5.45), Inches(10.5), Inches(0.55),
                    "— " + s.get("attribution", ""),
                    size=13, color=PURPLE_DARK, font=FONT_SANS, bold=True)

    def render_references(slide, s):
        add_slide_header(slide, s.get("title", "References"))
        items = s.get("items", [])[:14]
        if not items:
            return
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.85), Inches(12.13), Inches(5.3))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            n = it.get("n", "")
            src = it.get("source", "")
            ttl = it.get("title", "")
            url = it.get("url", "")
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            pPr = p._p
            for r in pPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
                pPr.remove(r)
            p.line_spacing = 1.22
            run_n = p.add_run()
            run_n.text = f"[{n}]  "
            run_n.font.name = FONT_SANS
            run_n.font.size = Pt(9.5)
            run_n.font.bold = True
            run_n.font.color.rgb = PURPLE
            rest = "  ".join(x for x in [src, ttl, url] if x)
            run_t = p.add_run()
            run_t.text = rest
            run_t.font.name = FONT_SANS
            run_t.font.size = Pt(9.5)
            run_t.font.color.rgb = TEXT_MUTED

    def render_closing(slide, s):
        add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
        # Right purple panel
        panel_x = Inches(9.2)
        add_rect(slide, panel_x, Emu(0), int(SLIDE_W) - int(panel_x), SLIDE_H, PURPLE)
        # Left content
        add_textbox(slide, Inches(0.8), Inches(2.4), Inches(8.0), Inches(2.0),
                    s.get("title", "Thank you"),
                    size=52, bold=True, color=PURPLE_DARK, font=FONT_SANS)
        add_rect(slide, Inches(0.8), Inches(4.5), Inches(2.8), Inches(0.12), PURPLE)
        if s.get("subtitle"):
            add_textbox(slide, Inches(0.8), Inches(4.75), Inches(8.0), Inches(0.7),
                        s["subtitle"], size=15, color=TEXT_MUTED, font=FONT_SANS)
        add_textbox(slide, Inches(0.8), Inches(6.75), Inches(8.0), Inches(0.4),
                    f"Generated by SkillCTI  ·  {date_str}  ·  {classification}",
                    size=9.5, color=TEXT_FAINT, font=FONT_SANS)
        # Branding on right panel
        add_textbox(slide, panel_x, Inches(3.2), int(SLIDE_W) - int(panel_x), Inches(0.65),
                    "SkillCTI", size=22, bold=True, color=WHITE,
                    font=FONT_SANS, align=PP_ALIGN.CENTER)
        add_textbox(slide, panel_x, Inches(3.95), int(SLIDE_W) - int(panel_x), Inches(0.45),
                    "CYBER THREAT INTELLIGENCE", size=8, bold=True, color=PURPLE_ACCENT,
                    font=FONT_SANS, align=PP_ALIGN.CENTER)

    # ── New visual slide types ─────────────────────────────────────

    def render_callout(slide, s):
        """2–4 severity-coded callout boxes — key findings, critical actions, risk items."""
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        callouts = s.get("callouts", [])[:4]
        if not callouts:
            return

        n = len(callouts)
        cols = 2 if n > 2 else n
        rows = (n + cols - 1) // cols

        gap = int(Inches(0.22))
        usable_w = int(Inches(12.13))
        usable_h = int(Inches(4.85))
        card_w = int((usable_w - gap * (cols - 1)) / cols)
        card_h = int((usable_h - gap * (rows - 1)) / rows)
        start_x = int(Inches(0.6))
        start_y = int(Inches(1.85))

        for i, c in enumerate(callouts):
            col = i % cols
            row = i // cols
            x = start_x + col * (card_w + gap)
            y = start_y + row * (card_h + gap)

            sev = c.get("severity", "neutral")
            fg, bg = sev_colors(sev)

            # Card base
            add_rect(slide, x, y, card_w, card_h, bg)
            # Thick colored top bar
            add_rect(slide, x, y, card_w, int(Inches(0.1)), fg)

            # Severity label badge (top right)
            if sev and sev not in ("neutral", ""):
                badge_w = int(Inches(1.0))
                add_textbox(slide,
                            x + card_w - badge_w - int(Inches(0.1)),
                            y + int(Inches(0.14)),
                            badge_w, int(Inches(0.28)),
                            sev.upper(), size=7.5, bold=True, color=fg,
                            font=FONT_SANS, align=PP_ALIGN.RIGHT)

            # Card title
            add_textbox(slide,
                        x + int(Inches(0.18)),
                        y + int(Inches(0.18)),
                        card_w - int(Inches(1.2)),
                        int(Inches(0.55)),
                        c.get("title", ""), size=13, bold=True, color=TEXT_DARK, font=FONT_SANS)

            # Divider under title
            add_rect(slide,
                     x + int(Inches(0.18)),
                     y + int(Inches(0.75)),
                     card_w - int(Inches(0.36)),
                     int(Emu(6000)), DIVIDER)

            # Body text
            add_textbox(slide,
                        x + int(Inches(0.18)),
                        y + int(Inches(0.88)),
                        card_w - int(Inches(0.36)),
                        card_h - int(Inches(1.0)),
                        c.get("body", ""), size=11.5, color=TEXT_DARK, font=FONT_SANS)

    def render_highlight(slide, s):
        """Large impact-statement slide — split layout: dark purple panel left, context right."""
        panel_w = int(Inches(5.6))
        add_rect(slide, Emu(0), Emu(0), panel_w, SLIDE_H, PURPLE_DARK)
        add_rect(slide, panel_w, Emu(0), int(SLIDE_W) - panel_w, SLIDE_H, WHITE)

        # Decorative subtle stripes on white area
        for i in range(6):
            add_rect(slide, panel_w, int(Inches(i * 1.25)), int(SLIDE_W) - panel_w, int(Emu(5000)), DIVIDER)

        # Classification label top right
        add_textbox(slide, int(Inches(9.2)), int(Inches(0.28)), int(Inches(3.9)), int(Inches(0.38)),
                    classification, size=8.5, color=TEXT_FAINT, font=FONT_SANS, align=PP_ALIGN.RIGHT)

        # Main value / stat on left (very large)
        value = s.get("value", s.get("message", ""))
        add_textbox(slide, int(Inches(0.3)), int(Inches(1.6)), panel_w - int(Inches(0.5)),
                    int(Inches(2.8)),
                    value, size=54, bold=True, color=WHITE, font=FONT_SANS,
                    anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

        # Label below value
        label = s.get("label", "")
        if label:
            add_rect(slide, int(Inches(0.5)), int(Inches(4.55)),
                     int(Inches(3.0)), int(Emu(7000)), PURPLE_ACCENT)
            add_textbox(slide, int(Inches(0.3)), int(Inches(4.75)),
                        panel_w - int(Inches(0.5)), int(Inches(0.5)),
                        label, size=12, color=PURPLE_LIGHT, font=FONT_SANS, align=PP_ALIGN.CENTER)

        # Right side: slide title
        right_x = panel_w + int(Inches(0.55))
        right_w = int(SLIDE_W) - panel_w - int(Inches(0.75))

        title = s.get("title", "")
        if title:
            add_textbox(slide, right_x, int(Inches(0.85)), right_w, int(Inches(0.65)),
                        title, size=12, bold=True, color=TEXT_MUTED, font=FONT_SANS)
            add_rect(slide, right_x, int(Inches(1.5)), right_w, int(Emu(7000)), DIVIDER)

        # Context paragraph
        context = s.get("context", "")
        if context:
            add_textbox(slide, right_x, int(Inches(1.75)), right_w, int(Inches(4.8)),
                        context, size=16, color=TEXT_DARK, font=FONT_SANS)

        # Bullets as fallback if no context text
        bullets = s.get("bullets", [])
        if bullets and not context:
            add_bullet_list(slide, bullets, right_x, int(Inches(1.75)), right_w,
                            int(Inches(4.8)), size=14)

    def render_timeline(slide, s):
        """Horizontal timeline — incident sequences, phase breakdowns, milestones."""
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        events = s.get("events", [])[:6]
        if not events:
            return

        n = len(events)
        line_y = int(Inches(4.0))
        line_h = int(Emu(9525 * 2))
        margin_l = int(Inches(0.9))
        margin_r = int(Inches(0.9))
        usable_w = int(SLIDE_W) - margin_l - margin_r

        # Main timeline line
        add_rect(slide, margin_l, line_y, usable_w, line_h, PURPLE_ACCENT)

        for i, ev in enumerate(events):
            if n > 1:
                cx = margin_l + int(usable_w * i / (n - 1))
            else:
                cx = margin_l + usable_w // 2

            sev = ev.get("severity", "neutral")
            dot_color, _ = sev_colors(sev)
            dot_r = int(Inches(0.14))
            dot_x = cx - dot_r
            dot_y = line_y - dot_r + line_h // 2

            # Event dot
            add_rect(slide, dot_x, dot_y, dot_r * 2, dot_r * 2, dot_color)

            label_w = int(Inches(1.9))
            label_x = cx - label_w // 2

            if i % 2 == 0:
                # Date above line
                add_textbox(slide, label_x, line_y - int(Inches(1.35)),
                            label_w, int(Inches(0.38)),
                            ev.get("date", ""), size=8.5, bold=True, color=PURPLE_DARK,
                            font=FONT_SANS, align=PP_ALIGN.CENTER)
                # Connector tick
                add_rect(slide, cx - int(Emu(3000)), line_y - int(Inches(0.55)),
                         int(Emu(6000)), int(Inches(0.55)), DIVIDER)
                # Label + detail below line
                add_textbox(slide, label_x, line_y + int(Inches(0.35)),
                            label_w, int(Inches(0.42)),
                            ev.get("label", ""), size=11, bold=True, color=TEXT_DARK,
                            font=FONT_SANS, align=PP_ALIGN.CENTER)
                if ev.get("detail"):
                    add_textbox(slide, label_x, line_y + int(Inches(0.78)),
                                label_w, int(Inches(0.72)),
                                ev["detail"], size=9, color=TEXT_MUTED,
                                font=FONT_SANS, align=PP_ALIGN.CENTER)
            else:
                # Label + detail above line
                if ev.get("detail"):
                    add_textbox(slide, label_x, line_y - int(Inches(1.65)),
                                label_w, int(Inches(0.72)),
                                ev["detail"], size=9, color=TEXT_MUTED,
                                font=FONT_SANS, align=PP_ALIGN.CENTER)
                add_textbox(slide, label_x, line_y - int(Inches(0.95)),
                            label_w, int(Inches(0.42)),
                            ev.get("label", ""), size=11, bold=True, color=TEXT_DARK,
                            font=FONT_SANS, align=PP_ALIGN.CENTER)
                # Connector tick
                add_rect(slide, cx - int(Emu(3000)), line_y + int(line_h),
                         int(Emu(6000)), int(Inches(0.55)), DIVIDER)
                # Date below line
                add_textbox(slide, label_x, line_y + int(Inches(0.6)),
                            label_w, int(Inches(0.38)),
                            ev.get("date", ""), size=8.5, bold=True, color=PURPLE_DARK,
                            font=FONT_SANS, align=PP_ALIGN.CENTER)

    def render_agenda(slide, s):
        """Numbered agenda / section overview slide."""
        add_slide_header(slide, s.get("title", "Agenda"), s.get("subtitle"))
        items = s.get("items", [])[:8]
        if not items:
            return

        n = len(items)
        cols = 2 if n > 4 else 1
        rows_per_col = (n + cols - 1) // cols
        item_h = int(Inches(0.75))
        start_y = int(Inches(1.72))
        col_w_total = int(Inches(12.13))
        col_w = int((col_w_total - int(Inches(0.4)) * (cols - 1)) / cols)
        left_margin = int(Inches(0.6))

        for i, item in enumerate(items):
            col = i // rows_per_col
            row = i % rows_per_col
            x = left_margin + col * (col_w + int(Inches(0.4)))
            y = start_y + row * item_h

            # Number badge
            badge = int(Inches(0.48))
            add_rect(slide, x, y + int(Inches(0.07)), badge, badge, PURPLE)
            add_textbox(slide, x, y + int(Inches(0.07)), badge, badge,
                        str(i + 1), size=14, bold=True, color=WHITE,
                        font=FONT_SANS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # Item label
            add_textbox(slide, x + badge + int(Inches(0.18)),
                        y + int(Inches(0.05)),
                        col_w - badge - int(Inches(0.25)),
                        int(Inches(0.58)),
                        str(item), size=15, color=TEXT_DARK, font=FONT_SANS,
                        anchor=MSO_ANCHOR.MIDDLE)

            # Row divider
            add_rect(slide, x, y + item_h - int(Emu(5000)),
                     col_w, int(Emu(5000)), DIVIDER)

    # ── Dispatch ───────────────────────────────────────────────────

    renderers = {
        "title":      render_title,
        "section":    render_section,
        "content":    render_content,
        "stats":      render_stats,
        "table":      render_table,
        "two_column": render_two_column,
        "quote":      render_quote,
        "references": render_references,
        "closing":    render_closing,
        "callout":    render_callout,
        "highlight":  render_highlight,
        "timeline":   render_timeline,
        "agenda":     render_agenda,
    }

    slides_data = outline.get("slides", []) if isinstance(outline, dict) else []
    if not slides_data:
        slides_data = [{"type": "title", "title": deck_title, "subtitle": "(empty deck)"}]

    total = len(slides_data)
    for idx, s in enumerate(slides_data, start=1):
        slide = add_blank_slide()
        bg = add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
        spTree = slide.shapes._spTree
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)
        stype = s.get("type", "content")
        renderer = renderers.get(stype, render_content)
        try:
            renderer(slide, s)
        except Exception as e:
            add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1.5),
                        f"[Slide {idx} render error: {type(e).__name__}: {e}]",
                        size=14, color=RGBColor(0xB9, 0x1C, 0x1C), font=FONT_SANS)
        if stype not in ("title", "section", "closing"):
            add_footer(slide, idx, total)

    prs.save(output_path)
