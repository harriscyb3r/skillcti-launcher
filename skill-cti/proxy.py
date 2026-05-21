"""
SkillCTI — Local API Proxy
---------------------------------
Keeps your Anthropic API key out of the browser entirely.
The HTML calls localhost:8765, this script forwards to api.anthropic.com.

Also stores generated reports on disk so the launcher can show a history view.
Reports are written next to this script in a ./reports/ directory:
  YYYY-MM-DDTHH-MM-SS-<skillId>.html       (open directly in any browser)
  YYYY-MM-DDTHH-MM-SS-<skillId>.meta.json  (skill metadata for the history list)

Usage:
  export ANTHROPIC_API_KEY=sk-ant-...
  python proxy.py
"""

import os
import json
import shutil
import subprocess
import tempfile
import urllib.request
import urllib.error
import urllib.parse
from pathlib import Path
from http.server import ThreadingHTTPServer, BaseHTTPRequestHandler

PORT = 8765
ANTHROPIC_API = "https://api.anthropic.com"
API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")

REPORTS_DIR = Path(__file__).resolve().parent / "reports"
REPORTS_DIR.mkdir(exist_ok=True)

# ── Headless-browser PDF conversion ───────────────────────────
# Uses Microsoft Edge or Google Chrome's headless mode to render HTML to a
# real vector PDF. Edge ships with Windows 11 by default. We cache the
# discovered browser path on first use.

_BROWSER_PATH = None

def find_browser():
    """Locate a Chromium-based browser for headless PDF rendering."""
    global _BROWSER_PATH
    if _BROWSER_PATH:
        return _BROWSER_PATH
    candidates = [
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    ]
    for name in ("msedge", "chrome", "google-chrome", "chromium", "chromium-browser"):
        p = shutil.which(name)
        if p:
            candidates.append(p)
    for c in candidates:
        if c and Path(c).exists():
            _BROWSER_PATH = c
            return c
    return None


def html_to_pdf(html_str, output_path):
    """Convert an HTML string to a PDF file using headless Edge/Chrome."""
    browser = find_browser()
    if not browser:
        raise RuntimeError(
            "No Chromium-based browser found for PDF conversion. "
            "Install Microsoft Edge (recommended on Windows) or Google Chrome."
        )

    # Write HTML to a temp file so the browser can load it as file:///
    with tempfile.NamedTemporaryFile(
        suffix=".html", delete=False, mode="w", encoding="utf-8"
    ) as f:
        f.write(html_str)
        html_path = Path(f.name)

    try:
        file_url = "file:///" + str(html_path).replace("\\", "/")
        cmd = [
            browser,
            "--headless=new",
            "--disable-gpu",
            "--no-sandbox",
            "--no-pdf-header-footer",
            "--run-all-compositor-stages-before-draw",
            "--virtual-time-budget=10000",
            f"--print-to-pdf={output_path}",
            file_url,
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=90)
        if result.returncode != 0 or not Path(output_path).exists():
            err = result.stderr.decode("utf-8", errors="replace")[-500:]
            raise RuntimeError(
                f"PDF conversion failed (exit {result.returncode}): {err.strip() or 'no stderr'}"
            )
    finally:
        try:
            html_path.unlink()
        except Exception:
            pass


# ── PowerPoint generator ──────────────────────────────────────
# Builds a professional white + deep-purple PPTX from a structured
# JSON slide outline produced by the model. See PPTX_OUTLINE_OVERRIDE
# in skills.js for the schema.

def outline_to_pptx(outline, output_path):
    """Render a slide-outline JSON object to a .pptx file.

    Theme: white background, deep purple primary, restrained accent use,
    Calibri sans-serif. Designed to look like a Mandiant / Bain analyst
    deck, not a generic Office template.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # Brand palette
    PURPLE      = RGBColor(0x6D, 0x28, 0xD9)  # deep, professional
    PURPLE_DARK = RGBColor(0x4C, 0x1D, 0x95)  # title text on title slides
    PURPLE_LIGHT= RGBColor(0xED, 0xE9, 0xFE)  # section-divider background tint
    PURPLE_ACCENT = RGBColor(0xA7, 0x8B, 0xFA)  # mid-purple for stats
    TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)  # near-black, warm
    TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)  # mid grey
    TEXT_FAINT  = RGBColor(0x9C, 0xA3, 0xAF)  # light grey
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    DIVIDER     = RGBColor(0xE5, 0xE7, 0xEB)

    FONT_SANS = "Calibri"  # universally available; renders cleanly

    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK_LAYOUT = prs.slide_layouts[6]  # blank

    meta = outline.get("meta", {}) if isinstance(outline, dict) else {}
    deck_title = meta.get("title", "Report")
    classification = meta.get("classification", "TLP:AMBER+STRICT")
    date_str = meta.get("date", "")
    author = meta.get("author", "SkillCTI")

    def add_blank_slide():
        return prs.slides.add_slide(BLANK_LAYOUT)

    def set_text(tf, text, *, size=12, bold=False, color=TEXT_DARK,
                 font=FONT_SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        # Drop any extra paragraphs beyond the first
        for p in list(tf.paragraphs)[1:]:
            tf._txBody.remove(p._p)
        # Defensive: ensure at least one paragraph exists (text frames
        # *should* always have one, but guard against malformed states)
        if not tf.paragraphs:
            tf.add_paragraph()
        p = tf.paragraphs[0]
        p.alignment = align
        # Strip any pre-existing runs cleanly and add a single fresh run
        pPr = p._p
        for r in pPr.findall(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}r'
        ):
            pPr.remove(r)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def add_textbox(slide, left, top, width, height, text, **kw):
        tb = slide.shapes.add_textbox(left, top, width, height)
        set_text(tb.text_frame, text, **kw)
        return tb

    def add_rect(slide, left, top, width, height, fill, line=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
        shp.shadow.inherit = False
        return shp

    def add_footer(slide, page_num, total):
        # Thin purple bar on the left edge, then footer text
        add_rect(slide, Emu(0), Inches(7.2), Inches(0.06), Inches(0.3), PURPLE)
        footer_text = f"{classification}  ·  {deck_title}  ·  {author}  ·  {page_num} / {total}"
        add_textbox(
            slide, Inches(0.2), Inches(7.15), Inches(13), Inches(0.35),
            footer_text, size=8, color=TEXT_FAINT, font=FONT_SANS,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

    def add_slide_header(slide, title, subtitle=None):
        # Purple accent bar
        add_rect(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.5), PURPLE)
        # Title
        add_textbox(
            slide, Inches(0.85), Inches(0.45), Inches(12), Inches(0.7),
            title, size=26, bold=True, color=TEXT_DARK, font=FONT_SANS,
        )
        if subtitle:
            add_textbox(
                slide, Inches(0.85), Inches(1.1), Inches(12), Inches(0.4),
                subtitle, size=12, color=TEXT_MUTED, font=FONT_SANS,
            )
        # Horizontal rule under the title block
        add_rect(slide, Inches(0.6), Inches(1.6), Inches(12.15), Emu(9525), DIVIDER)

    def style_run(run, text, *, size=14, bold=False, color=TEXT_DARK, font=FONT_SANS):
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def fill_bullet_paragraph(p, item, size, line_spacing, bullet_color):
        """Populate an existing paragraph with bullet glyph + lead-bold + body."""
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Remove any pre-existing runs (e.g. the default empty run on a fresh
        # paragraph) so we control the formatting cleanly.
        pPr = p._p
        for r in pPr.findall(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}r'
        ):
            pPr.remove(r)
        # Bullet glyph in the accent colour
        style_run(p.add_run(), "•  ", size=size, bold=True, color=bullet_color)
        # Optional "Lead phrase:" bold detection — only on short leads
        text = str(item)
        if ":" in text and text.index(":") < 40:
            head, tail = text.split(":", 1)
            style_run(p.add_run(), head + ":", size=size, bold=True, color=TEXT_DARK)
            style_run(p.add_run(), tail, size=size, color=TEXT_DARK)
        else:
            style_run(p.add_run(), text, size=size, color=TEXT_DARK)

    def add_bullet_list(slide, bullets, left, top, width, height,
                        size=14, line_spacing=1.35, bullet_color=PURPLE):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        # A fresh text frame always has exactly one default paragraph. Reuse
        # it for the first bullet, then add subsequent paragraphs. Never
        # remove all paragraphs — OOXML requires at least one <a:p>, and
        # accessing tf.paragraphs[0] on an empty frame raises IndexError.
        if not bullets:
            return
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            fill_bullet_paragraph(p, item, size, line_spacing, bullet_color)

    # ── Render each slide type ────────────────────────────────

    def render_title(slide, s):
        # White background with subtle purple sidebar on the left
        add_rect(slide, Emu(0), Emu(0), Inches(0.8), SLIDE_H, PURPLE)
        # Classification pill, top-right
        pill_w, pill_h = Inches(2.0), Inches(0.35)
        pill = add_rect(slide, Inches(11.1), Inches(0.5), pill_w, pill_h, WHITE, line=PURPLE)
        add_textbox(
            slide, Inches(11.1), Inches(0.5), pill_w, pill_h,
            classification, size=10, bold=True, color=PURPLE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # Main title — large, weighted, deep purple
        add_textbox(
            slide, Inches(1.2), Inches(2.6), Inches(11.5), Inches(1.6),
            s.get("title", deck_title),
            size=44, bold=True, color=PURPLE_DARK, font=FONT_SANS,
        )
        # Subtitle
        sub = s.get("subtitle") or meta.get("subtitle", "")
        if sub:
            add_textbox(
                slide, Inches(1.2), Inches(4.2), Inches(11.5), Inches(0.6),
                sub, size=18, color=TEXT_MUTED, font=FONT_SANS,
            )
        # Date + author at bottom
        client_line = meta.get("client", "")
        bottom = "  ·  ".join([x for x in [client_line, date_str, author] if x])
        add_textbox(
            slide, Inches(1.2), Inches(6.6), Inches(11.5), Inches(0.4),
            bottom, size=11, color=TEXT_FAINT, font=FONT_SANS,
        )

    def render_section(slide, s):
        # Full purple-tinted background
        add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, PURPLE_LIGHT)
        # Vertical purple stripe on left
        add_rect(slide, Inches(0.6), Inches(2.5), Inches(0.15), Inches(2.5), PURPLE)
        add_textbox(
            slide, Inches(1.1), Inches(2.8), Inches(11.5), Inches(1.2),
            s.get("title", "Section"),
            size=40, bold=True, color=PURPLE_DARK, font=FONT_SANS,
        )
        if s.get("subtitle"):
            add_textbox(
                slide, Inches(1.1), Inches(4.2), Inches(11.5), Inches(0.6),
                s["subtitle"], size=16, color=TEXT_MUTED, font=FONT_SANS,
            )

    def render_content(slide, s):
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        bullets = s.get("bullets", [])
        if bullets:
            add_bullet_list(
                slide, bullets,
                Inches(0.85), Inches(1.9), Inches(12), Inches(5.0),
                size=15,
            )
        if s.get("footnote"):
            add_textbox(
                slide, Inches(0.85), Inches(6.7), Inches(12), Inches(0.4),
                s["footnote"], size=9, color=TEXT_FAINT, font=FONT_SANS,
            )

    def render_stats(slide, s):
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        stats = s.get("stats", [])[:5]
        if not stats:
            return
        n = len(stats)
        gap = Inches(0.3)
        usable_w = Inches(12.0)
        card_w = Emu((usable_w - gap * (n - 1)) // n)
        left = Inches(0.85)
        top = Inches(2.5)
        card_h = Inches(2.5)
        for i, st in enumerate(stats):
            x = Emu(left + i * (card_w + gap))
            # subtle purple-tinted card
            add_rect(slide, x, top, card_w, card_h, WHITE, line=DIVIDER)
            # accent stripe at the top of the card
            add_rect(slide, x, top, card_w, Inches(0.08), PURPLE)
            # Big value
            add_textbox(
                slide, x, Emu(top + Inches(0.5)), card_w, Inches(1.1),
                str(st.get("value", "")),
                size=36, bold=True, color=PURPLE_DARK, font=FONT_SANS,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            # Label
            add_textbox(
                slide, x, Emu(top + Inches(1.65)), card_w, Inches(0.7),
                str(st.get("label", "")),
                size=11, color=TEXT_MUTED, font=FONT_SANS,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
            )

    def render_table(slide, s):
        add_slide_header(slide, s.get("title", ""), s.get("subtitle"))
        headers = s.get("headers", [])
        rows = s.get("rows", [])
        if not headers or not rows:
            return
        n_cols = len(headers)
        n_rows = len(rows) + 1
        table = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.85), Inches(1.95),
            Inches(12), min(Inches(0.5) * n_rows, Inches(5.0)),
        ).table
        # Header row
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PURPLE
            set_text(cell.text_frame, str(h), size=11, bold=True,
                     color=WHITE, font=FONT_SANS, align=PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
        # Body rows
        for i, row in enumerate(rows, start=1):
            for j in range(n_cols):
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF9, 0xFA, 0xFB)
                value = str(row[j]) if j < len(row) else ""
                set_text(cell.text_frame, value, size=10, color=TEXT_DARK,
                         font=FONT_SANS, align=PP_ALIGN.LEFT,
                         anchor=MSO_ANCHOR.MIDDLE)

    def render_two_column(slide, s):
        add_slide_header(slide, s.get("title", ""))
        col_w = Inches(5.8)
        gap = Inches(0.4)
        left_x = Inches(0.85)
        right_x = Emu(left_x + col_w + gap)
        top = Inches(2.0)
        # Left column header
        add_textbox(
            slide, left_x, top, col_w, Inches(0.5),
            s.get("left_title", "Left"),
            size=14, bold=True, color=PURPLE_DARK, font=FONT_SANS,
        )
        add_rect(slide, left_x, Emu(top + Inches(0.5)), col_w, Emu(9525), PURPLE)
        add_bullet_list(
            slide, s.get("left_bullets", []),
            left_x, Emu(top + Inches(0.7)), col_w, Inches(4.5),
            size=13,
        )
        # Right column header
        add_textbox(
            slide, right_x, top, col_w, Inches(0.5),
            s.get("right_title", "Right"),
            size=14, bold=True, color=PURPLE_DARK, font=FONT_SANS,
        )
        add_rect(slide, right_x, Emu(top + Inches(0.5)), col_w, Emu(9525), PURPLE)
        add_bullet_list(
            slide, s.get("right_bullets", []),
            right_x, Emu(top + Inches(0.7)), col_w, Inches(4.5),
            size=13,
        )

    def render_quote(slide, s):
        # Subtle purple tint background
        add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, RGBColor(0xFB, 0xFA, 0xFE))
        # Big opening quote mark in purple
        add_textbox(
            slide, Inches(1.0), Inches(1.4), Inches(2), Inches(2),
            "“", size=140, bold=True, color=PURPLE_ACCENT, font=FONT_SANS,
        )
        # Quote text
        add_textbox(
            slide, Inches(2.0), Inches(2.5), Inches(10), Inches(2.5),
            s.get("quote", ""),
            size=24, color=TEXT_DARK, font=FONT_SANS,
        )
        # Attribution
        add_textbox(
            slide, Inches(2.0), Inches(5.5), Inches(10), Inches(0.5),
            "— " + s.get("attribution", ""),
            size=13, color=PURPLE_DARK, font=FONT_SANS, bold=True,
        )

    def render_references(slide, s):
        add_slide_header(slide, s.get("title", "References"))
        items = s.get("items", [])[:12]
        lines = []
        for it in items:
            n = it.get("n", "")
            src = it.get("source", "")
            ttl = it.get("title", "")
            url = it.get("url", "")
            lines.append(f"[{n}] {src} — {ttl} — {url}".strip(" —"))
        if lines:
            tb = slide.shapes.add_textbox(
                Inches(0.85), Inches(1.9), Inches(12), Inches(5.2)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            # Reuse the default first paragraph for line 0; add new ones after
            for i, ln in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                # Strip any pre-existing runs on the reused paragraph
                pPr = p._p
                for r in pPr.findall(
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}r'
                ):
                    pPr.remove(r)
                p.line_spacing = 1.3
                run = p.add_run()
                run.text = ln
                run.font.name = FONT_SANS
                run.font.size = Pt(10)
                run.font.color.rgb = TEXT_MUTED

    def render_closing(slide, s):
        add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
        # Big purple band on left
        add_rect(slide, Emu(0), Emu(0), Inches(0.8), SLIDE_H, PURPLE)
        add_textbox(
            slide, Inches(1.2), Inches(3.0), Inches(11.5), Inches(1.5),
            s.get("title", "Thank you"),
            size=52, bold=True, color=PURPLE_DARK, font=FONT_SANS,
            align=PP_ALIGN.LEFT,
        )
        if s.get("subtitle"):
            add_textbox(
                slide, Inches(1.2), Inches(4.5), Inches(11.5), Inches(0.6),
                s["subtitle"], size=18, color=TEXT_MUTED, font=FONT_SANS,
            )
        # Footer line
        add_textbox(
            slide, Inches(1.2), Inches(6.7), Inches(11.5), Inches(0.4),
            f"Generated by SkillCTI  ·  {date_str}  ·  {classification}",
            size=10, color=TEXT_FAINT, font=FONT_SANS,
        )

    renderers = {
        "title":      render_title,
        "section":    render_section,
        "content":    render_content,
        "stats":      render_stats,
        "table":      render_table,
        "two_column": render_two_column,
        "quote":      render_quote,
        "references": render_references,
        "closing":    render_closing,
    }

    slides = outline.get("slides", []) if isinstance(outline, dict) else []
    if not slides:
        # Defensive: create a minimal title slide so file isn't empty
        slides = [{"type": "title", "title": deck_title, "subtitle": "(empty deck)"}]

    total = len(slides)
    for idx, s in enumerate(slides, start=1):
        slide = add_blank_slide()
        # White background for all non-tinted slides
        bg = add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
        # Move background to back
        spTree = slide.shapes._spTree
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)
        stype = s.get("type", "content")
        renderer = renderers.get(stype, render_content)
        try:
            renderer(slide, s)
        except Exception as e:
            # Don't fail the whole deck — render an error placeholder
            add_textbox(
                slide, Inches(1), Inches(3), Inches(11), Inches(1.5),
                f"[Slide {idx} render error: {type(e).__name__}: {e}]",
                size=14, color=RGBColor(0xB9, 0x1C, 0x1C), font=FONT_SANS,
            )
        # Footer on every slide except title and closing
        if stype not in ("title", "closing"):
            add_footer(slide, idx, total)

    prs.save(output_path)


def safe_id(s):
    """URL-decode and reject path traversal. Preserves spaces, dots, parens —
    anything that's a legal Windows/Unix filename character."""
    decoded = urllib.parse.unquote(s or "")
    # Take just the basename, blocking both / and \ as path separators
    decoded = decoded.replace("\\", "/").split("/")[-1]
    # Block path traversal and empty
    if not decoded or decoded in (".", ".."):
        return ""
    # Strip characters that are dangerous in JS string contexts or filenames
    decoded = "".join(c for c in decoded if c not in "<>\"'\x00\r\n")
    return decoded[:200]


def adopt_orphan_html(html_file):
    """Create a meta file for an HTML report dropped into reports/ by hand.
    If the filename contains characters that would break URL or JS round-tripping,
    rename the file to a safe stem first so the meta id and html filename match."""
    from datetime import datetime, timezone
    original_stem = html_file.stem
    # Replace anything that would break JS string contexts or URLs
    safe_stem = "".join(c if c not in "<>\"'\x00\r\n" else "_" for c in original_stem)
    if safe_stem != original_stem:
        new_path = REPORTS_DIR / f"{safe_stem}.html"
        html_file.rename(new_path)
        html_file = new_path
    stem = safe_stem
    mtime = datetime.fromtimestamp(html_file.stat().st_mtime, tz=timezone.utc)
    title = stem.replace("-", " ").replace("_", " ").strip() or stem
    meta = {
        "id": stem,
        "skillId": "imported",
        "skillName": "Imported Report",
        "badge": "IMPORTED",
        "badgeColor": "#6b7280",
        "timestamp": mtime.isoformat(),
        "title": title[:80],
        "inputs": {},
        "bytes": html_file.stat().st_size,
    }
    (REPORTS_DIR / f"{stem}.meta.json").write_text(
        json.dumps(meta, indent=2), encoding="utf-8"
    )
    return meta


def list_reports():
    """Return all saved reports' metadata, newest first.
    Auto-adopts any .html file in reports/ that doesn't have a paired .meta.json."""
    for html_file in REPORTS_DIR.glob("*.html"):
        meta_file = html_file.with_suffix(".meta.json")
        # .with_suffix only replaces the last suffix, so handle the .meta.json case explicitly
        meta_file = REPORTS_DIR / f"{html_file.stem}.meta.json"
        if not meta_file.exists():
            try:
                adopt_orphan_html(html_file)
                print(f"  · adopted orphan report: {html_file.name}")
            except Exception as e:
                print(f"  · failed to adopt {html_file.name}: {e}")

    items = []
    for meta_file in REPORTS_DIR.glob("*.meta.json"):
        try:
            meta = json.loads(meta_file.read_text(encoding="utf-8"))
            mid = meta.get("id", "")
            # Skip meta files whose content partner has been deleted (any
            # of .html / .pdf / .pptx must still be on disk)
            html_ok = (REPORTS_DIR / f"{mid}.html").exists()
            pdf_ok  = (REPORTS_DIR / f"{mid}.pdf").exists()
            pptx_ok = (REPORTS_DIR / f"{mid}.pptx").exists()
            if not (html_ok or pdf_ok or pptx_ok):
                continue
            items.append(meta)
        except Exception:
            continue
    items.sort(key=lambda m: m.get("timestamp", ""), reverse=True)
    return items


class ProxyHandler(BaseHTTPRequestHandler):

    def log_message(self, format, *args):
        print(f"  → {args[0]} {args[1]}")

    def handle_one_request(self):
        # Swallow the noisy WinError 10053 / BrokenPipe / ConnectionReset that
        # happens when the browser aborts a request mid-flight (e.g. a 3s health
        # check timeout firing before the upstream API responds). These are
        # normal client-disconnect events, not bugs.
        try:
            super().handle_one_request()
        except (ConnectionAbortedError, BrokenPipeError, ConnectionResetError) as e:
            print(f"  · client disconnected ({type(e).__name__})")

    def send_cors(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, DELETE, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")

    def send_json(self, status, payload):
        self.send_response(status)
        self.send_cors()
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        self.wfile.write(json.dumps(payload).encode("utf-8"))

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_cors()
        self.end_headers()

    # ── GET — list reports or fetch a single one ──────────────

    def do_GET(self):
        if self.path == "/health":
            self.send_json(200, {"ok": True})
            return

        if self.path == "/reports":
            self.send_json(200, {"reports": list_reports()})
            return

        # Ransomware.live passthrough for the dashboard widget — fetches the
        # AU country victims feed server-side to avoid CORS issues in the
        # browser. Response is normalised to a small subset for the widget.
        #
        # Note: the country code MUST be lowercase. Uppercase "AU" returns
        # 404; lowercase "au" works against the documented v2 endpoint.
        # Verified 2026-05-21 against the apidocs at
        # https://www.ransomware.live/apidocs.
        if self.path == "/ransomware-au":
            RANSOMWARE_AU_URL = "https://api.ransomware.live/v2/countryvictims/au"
            try:
                req = urllib.request.Request(
                    RANSOMWARE_AU_URL,
                    headers={
                        "User-Agent": "skillcti-dashboard/1.0",
                        "Accept": "application/json",
                    },
                )
                with urllib.request.urlopen(req, timeout=10) as r:
                    raw = r.read().decode("utf-8")
                data = json.loads(raw)
                # The v2 API returns either a list or {"victims":[...]}. Normalise.
                victims = data if isinstance(data, list) else data.get("victims", []) or data.get("data", [])
                self.send_json(200, {"victims": victims})
            except Exception as e:
                self.send_json(502, {"error": f"ransomware.live fetch failed: {e}"})
            return

        # /reports/<id>/pdf — return the PDF binary directly (for in-tab viewing)
        if self.path.startswith("/reports/") and self.path.endswith("/pdf"):
            report_id = safe_id(self.path[len("/reports/"):-len("/pdf")])
            pdf_file = REPORTS_DIR / f"{report_id}.pdf"
            if not pdf_file.exists():
                self.send_json(404, {"error": "pdf not found"})
                return
            pdf_bytes = pdf_file.read_bytes()
            self.send_response(200)
            self.send_cors()
            self.send_header("Content-Type", "application/pdf")
            self.send_header("Content-Disposition", f'inline; filename="{report_id}.pdf"')
            self.send_header("Content-Length", str(len(pdf_bytes)))
            self.end_headers()
            self.wfile.write(pdf_bytes)
            return

        # /reports/<id>/pptx — return the PowerPoint binary directly
        if self.path.startswith("/reports/") and self.path.endswith("/pptx"):
            report_id = safe_id(self.path[len("/reports/"):-len("/pptx")])
            pptx_file = REPORTS_DIR / f"{report_id}.pptx"
            if not pptx_file.exists():
                self.send_json(404, {"error": "pptx not found"})
                return
            pptx_bytes = pptx_file.read_bytes()
            self.send_response(200)
            self.send_cors()
            self.send_header(
                "Content-Type",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            self.send_header("Content-Disposition", f'attachment; filename="{report_id}.pptx"')
            self.send_header("Content-Length", str(len(pptx_bytes)))
            self.end_headers()
            self.wfile.write(pptx_bytes)
            return

        if self.path.startswith("/reports/"):
            report_id = safe_id(self.path[len("/reports/"):])
            html_file = REPORTS_DIR / f"{report_id}.html"
            meta_file = REPORTS_DIR / f"{report_id}.meta.json"
            if not meta_file.exists():
                self.send_json(404, {"error": "report not found"})
                return
            meta = json.loads(meta_file.read_text(encoding="utf-8"))
            # PDF reports — there's no HTML to send back, just the meta
            if meta.get("format") == "pdf" or not html_file.exists():
                self.send_json(200, {"meta": meta, "html": None})
                return
            html = html_file.read_text(encoding="utf-8")
            self.send_json(200, {"meta": meta, "html": html})
            return

        self.send_json(404, {"error": "not found"})

    # ── DELETE — remove a report ──────────────────────────────

    def do_DELETE(self):
        if self.path.startswith("/reports/"):
            report_id = safe_id(self.path[len("/reports/"):])
            removed = False
            for ext in (".html", ".pdf", ".pptx", ".meta.json"):
                p = REPORTS_DIR / f"{report_id}{ext}"
                if p.exists():
                    p.unlink()
                    removed = True
            if removed:
                self.send_json(200, {"ok": True})
            else:
                self.send_json(404, {"error": "report not found"})
            return
        self.send_json(404, {"error": "not found"})

    # ── POST — save a report, or forward to Anthropic ─────────

    def do_POST(self):
        length = int(self.headers.get("Content-Length", 0))
        body = self.rfile.read(length)

        # PDF conversion endpoint — convert HTML → real .pdf via headless Edge,
        # save to /reports/<id>.pdf + meta, stream the PDF binary back so the
        # browser triggers a download.
        if self.path == "/generate-pdf":
            try:
                payload = json.loads(body.decode("utf-8"))
                meta = payload.get("meta") or {}
                html = payload.get("html") or ""
                if not html or not meta.get("id"):
                    self.send_json(400, {"error": "missing html or meta.id"})
                    return
                report_id = safe_id(meta["id"])
                meta["id"] = report_id
                meta["format"] = "pdf"
                pdf_path = REPORTS_DIR / f"{report_id}.pdf"
                html_to_pdf(html, str(pdf_path))
                meta["bytes"] = pdf_path.stat().st_size
                (REPORTS_DIR / f"{report_id}.meta.json").write_text(
                    json.dumps(meta, indent=2), encoding="utf-8"
                )
                pdf_bytes = pdf_path.read_bytes()
                self.send_response(200)
                self.send_cors()
                self.send_header("Content-Type", "application/pdf")
                self.send_header(
                    "Content-Disposition",
                    f'attachment; filename="{report_id}.pdf"',
                )
                self.send_header("Content-Length", str(len(pdf_bytes)))
                self.send_header("X-Report-Id", report_id)
                self.end_headers()
                self.wfile.write(pdf_bytes)
            except Exception as e:
                self.send_json(500, {"error": str(e)})
            return

        # PPTX generation endpoint — accepts a slide-outline JSON object
        # from the model and renders it to a professional white + purple
        # .pptx via python-pptx. Saves to /reports/<id>.pptx + meta and
        # streams the binary back so the browser triggers a download.
        if self.path == "/generate-pptx":
            try:
                payload = json.loads(body.decode("utf-8"))
                meta = payload.get("meta") or {}
                outline = payload.get("outline") or {}
                if not outline or not meta.get("id"):
                    self.send_json(400, {"error": "missing outline or meta.id"})
                    return
                report_id = safe_id(meta["id"])
                meta["id"] = report_id
                meta["format"] = "pptx"
                pptx_path = REPORTS_DIR / f"{report_id}.pptx"
                outline_to_pptx(outline, str(pptx_path))
                meta["bytes"] = pptx_path.stat().st_size
                (REPORTS_DIR / f"{report_id}.meta.json").write_text(
                    json.dumps(meta, indent=2), encoding="utf-8"
                )
                pptx_bytes = pptx_path.read_bytes()
                self.send_response(200)
                self.send_cors()
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition",
                    f'attachment; filename="{report_id}.pptx"',
                )
                self.send_header("Content-Length", str(len(pptx_bytes)))
                self.send_header("X-Report-Id", report_id)
                self.end_headers()
                self.wfile.write(pptx_bytes)
            except Exception as e:
                import traceback
                tb = traceback.format_exc()[-1200:]
                print(f"  ✗ PPTX generation error: {e}\n{tb}")
                self.send_json(500, {"error": f"PPTX generation failed: {e}"})
            return

        # Save endpoint — handled locally
        if self.path == "/reports":
            try:
                payload = json.loads(body.decode("utf-8"))
                meta = payload.get("meta") or {}
                html = payload.get("html") or ""
                if not html or not meta.get("id"):
                    self.send_json(400, {"error": "missing html or meta.id"})
                    return
                report_id = safe_id(meta["id"])
                meta["id"] = report_id  # normalised
                (REPORTS_DIR / f"{report_id}.html").write_text(html, encoding="utf-8")
                (REPORTS_DIR / f"{report_id}.meta.json").write_text(
                    json.dumps(meta, indent=2), encoding="utf-8"
                )
                self.send_json(200, {"ok": True, "id": report_id})
            except Exception as e:
                self.send_json(500, {"error": str(e)})
            return

        # Anything else — forward to Anthropic
        if not API_KEY:
            self.send_json(500, {
                "error": "ANTHROPIC_API_KEY environment variable is not set."
            })
            return

        target_url = ANTHROPIC_API + self.path
        req = urllib.request.Request(
            target_url,
            data=body,
            headers={
                "Content-Type": "application/json",
                "x-api-key": API_KEY,
                "anthropic-version": "2023-06-01",
                "anthropic-beta": "web-search-2025-03-05",
            },
            method="POST",
        )

        try:
            with urllib.request.urlopen(req) as resp:
                resp_body = resp.read()
                self.send_response(resp.status)
                self.send_cors()
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(resp_body)
                # Print prompt-cache + token usage stats to the proxy terminal
                # so you can see caching working in real time.
                try:
                    parsed = json.loads(resp_body)
                    usage = parsed.get("usage") or {}
                    if usage:
                        fresh   = usage.get("input_tokens", 0)
                        written = usage.get("cache_creation_input_tokens", 0)
                        read    = usage.get("cache_read_input_tokens", 0)
                        out     = usage.get("output_tokens", 0)
                        total_in = fresh + written + read
                        pct = (read / total_in * 100) if total_in else 0
                        tag = (
                            f"CACHE MISS · wrote {written:,} tokens"   if written and not read
                            else f"CACHE HIT  · {pct:.0f}% ({read:,}/{total_in:,})" if read
                            else "no cache"
                        )
                        print(f"  · usage: in={total_in:,} out={out:,} · {tag}")
                except Exception:
                    pass
        except urllib.error.HTTPError as e:
            err_body = e.read()
            self.send_response(e.code)
            self.send_cors()
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(err_body)
        except Exception as e:
            self.send_json(502, {"error": str(e)})


if __name__ == "__main__":
    if not API_KEY:
        print("\n  ✗  ANTHROPIC_API_KEY is not set.")
        print("     Run:  export ANTHROPIC_API_KEY=sk-ant-...")
        print("     Then: python proxy.py\n")
        exit(1)

    masked = API_KEY[:12] + "..." + API_KEY[-4:]
    browser = find_browser()
    browser_status = browser if browser else "NOT FOUND — PDF mode will fail"
    print(f"\n  SkillCTI Proxy")
    print(f"  ─────────────────────────────────")
    print(f"  API key  : {masked}")
    print(f"  Reports  : {REPORTS_DIR}")
    print(f"  Browser  : {browser_status}")
    print(f"  Listening on http://localhost:{PORT}")
    print(f"  Open skill-cti.html in Chrome")
    print(f"  Press Ctrl+C to stop\n")

    server = ThreadingHTTPServer(("localhost", PORT), ProxyHandler)
    server.daemon_threads = True  # so Ctrl+C exits cleanly even mid-request
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n  Proxy stopped.\n")
